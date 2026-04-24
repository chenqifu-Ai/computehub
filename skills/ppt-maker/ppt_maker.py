#!/usr/bin/env python3
"""
PPT制作工具 - 将Markdown内容转换为PowerPoint演示文稿
"""

import argparse
import json
import re
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE
except ImportError:
    print("请先安装依赖: pip install python-pptx")
    exit(1)

class PPTMaker:
    def __init__(self, template="business"):
        """初始化PPT制作器"""
        self.presentation = Presentation()
        self.template = template
        self.setup_template()
        
    def setup_template(self):
        """设置模板样式"""
        # 默认使用空白模板，可以根据需要扩展
        self.title_size = Pt(44)
        self.content_size = Pt(28)
        self.subtitle_size = Pt(32)
        
        # 颜色配置
        if self.template == "business":
            self.primary_color = RGBColor(0, 120, 215)  # 蓝色
            self.secondary_color = RGBColor(79, 129, 189)  # 浅蓝
            self.text_color = RGBColor(0, 0, 0)  # 黑色
            self.background_color = RGBColor(255, 255, 255)  # 白色
        elif self.template == "creative":
            self.primary_color = RGBColor(255, 107, 107)  # 红色
            self.secondary_color = RGBColor(255, 159, 67)  # 橙色
            self.text_color = RGBColor(51, 51, 51)  # 深灰
            self.background_color = RGBColor(249, 249, 249)  # 浅灰
        else:
            # 默认样式
            self.primary_color = RGBColor(0, 120, 215)
            self.secondary_color = RGBColor(79, 129, 189)
            self.text_color = RGBColor(0, 0, 0)
            self.background_color = RGBColor(255, 255, 255)
    
    def add_title_slide(self, title, subtitle=""):
        """添加标题幻灯片"""
        slide_layout = self.presentation.slide_layouts[0]  # 标题幻灯片
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # 设置样式
        self._set_text_style(title_shape, self.title_size, self.primary_color)
        self._set_text_style(subtitle_shape, self.subtitle_size, self.secondary_color)
        
        return slide
    
    def add_content_slide(self, title, content):
        """添加内容幻灯片"""
        slide_layout = self.presentation.slide_layouts[1]  # 标题和内容
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
        
        # 设置样式
        self._set_text_style(title_shape, self.title_size, self.primary_color)
        self._set_text_style(content_shape, self.content_size, self.text_color)
        
        return slide
    
    def add_bullet_slide(self, title, items):
        """添加项目符号幻灯片"""
        slide = self.add_content_slide(title, "")
        content_shape = slide.placeholders[1]
        
        # 清空原有内容
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # 添加项目符号
        for item in items:
            p = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 0 else text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = self.content_size
            p.font.color.rgb = self.text_color
        
        return slide
    
    def _set_text_style(self, shape, font_size, color):
        """设置文本样式"""
        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.color.rgb = color
                run.font.name = "Arial"
    
    def save(self, filename):
        """保存PPT文件"""
        self.presentation.save(filename)
        print(f"✅ PPT文件已保存: {filename}")

def markdown_to_ppt(md_content, output_file, template="business"):
    """将Markdown内容转换为PPT"""
    maker = PPTMaker(template)
    
    # 解析Markdown内容
    lines = md_content.split('\n')
    current_slide_title = None
    current_content = []
    
    for line in lines:
        # 检测幻灯片标题 (## 开头)
        if line.startswith('## '):
            # 保存前一页幻灯片
            if current_slide_title and current_content:
                maker.add_bullet_slide(current_slide_title, current_content)
            
            # 开始新幻灯片
            current_slide_title = line[3:].strip()  # 移除 '## '
            current_content = []
        
        # 检测项目符号
        elif line.startswith('- ') or line.startswith('* '):
            current_content.append(line[2:].strip())
        
        # 检测普通文本
        elif line.strip() and not line.startswith('#') and not line.startswith('---'):
            current_content.append(line.strip())
    
    # 添加最后一页幻灯片
    if current_slide_title and current_content:
        maker.add_bullet_slide(current_slide_title, current_content)
    
    # 保存文件
    maker.save(output_file)
    return output_file

def convert_markdown_file(md_file, output_file, template="business"):
    """转换Markdown文件到PPT"""
    try:
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return markdown_to_ppt(content, output_file, template)
    except Exception as e:
        print(f"❌ 转换失败: {e}")
        return None

def main():
    parser = argparse.ArgumentParser(description="Markdown转PPT工具")
    parser.add_argument("input", help="输入Markdown文件")
    parser.add_argument("output", help="输出PPT文件")
    parser.add_argument("--template", "-t", default="business", 
                       choices=["business", "creative", "technical"],
                       help="模板风格")
    
    args = parser.parse_args()
    
    # 检查输入文件
    if not Path(args.input).exists():
        print(f"❌ 输入文件不存在: {args.input}")
        return
    
    # 转换文件
    result = convert_markdown_file(args.input, args.output, args.template)
    
    if result:
        print(f"🎉 转换成功!")
        print(f"📁 输入: {args.input}")
        print(f"📁 输出: {args.output}")
        print(f"🎨 模板: {args.template}")
    else:
        print("❌ 转换失败")

if __name__ == "__main__":
    main()