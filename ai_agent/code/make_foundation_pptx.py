#!/usr/bin/env python3
"""从PDF提取的内容生成世界华人基金会简介PPTX，提交ECSS视频管线"""
import json, os, sys, subprocess, tempfile

slides = [
    {
        "title": "世界华人基金会",
        "body": "凝聚华人力量 · 弘扬中华文明 · 推动交流合作\n\n宗旨理念：和平、友爱、发展、共赢",
        "sub": True
    },
    {
        "title": "核心宗旨",
        "body": "凝聚世界5000多万海外华人华侨的力量\n让全球华人紧密团结在一起"
    },
    {
        "title": "弘扬中华文明",
        "body": "通过各类活动弘扬中华文明\n传播中华文化的魅力\n增强华人对民族文化的认同感"
    },
    {
        "title": "推动交流合作",
        "body": "积极推动华人华侨与中国大陆之间的交流与合作\n促进双方在经济、文化等多领域共同发展"
    },
    {
        "title": "十年成就概览",
        "body": "• 华商经济合作 数百亿美元\n• 在大陆兴办外资合资企业 几千家\n• 组织全球商贸投资活动 数百次\n• 创办《世界杰出华人奖》\n• 促进中国与国外交流合作"
    },
    {
        "title": "主要贡献 · 经济合作",
        "body": "成立以来，促成世界各国华商经济合作数百亿美元\n为全球华商搭建了广阔的合作平台\n通过牵线搭桥，带动大量投资"
    },
    {
        "title": "主要贡献 · 兴办企业",
        "body": "到大陆内地成功兴办的外资、合资企业达几千家\n有力推动了中国大陆的经济发展\n覆盖多个产业领域"
    },
    {
        "title": "主要贡献 · 商贸活动",
        "body": "组织全球性商贸投资活动数百次\n包括：\n• 海峡两岸经贸洽谈会\n• 世界华商领袖年会\n• 世界华人小姐选美大赛等"
    },
    {
        "title": "主要贡献 · 杰出华人奖",
        "body": "创办及组织了多届《世界杰出华人奖》\n激励华人在各领域积极进取\n提升华人在全球的影响力"
    },
    {
        "title": "主要贡献 · 国际合作",
        "body": "促进中国与国外的交流合作\n对\"一带一路\"国家投资积极\n为国际合作与发展贡献力量"
    },
    {
        "title": "启动仪式 · 海南博鳌",
        "body": "世界华人协会红色文化展揭幕仪式\n暨世界华人名人奖启动仪式\n在海南省博鳌镇圆满收官\n主题：\"百年侨史、赤子侨心\""
    },
    {
        "title": "现有布局",
        "body": "已设立办事处：北京 · 福建 · 浙江 · 广东\n投资各类项目15个\n总价值数十亿元人民币"
    },
    {
        "title": "未来规划",
        "body": "• 今年内设立20个省级办事处和分会\n• 对接各地政府项目投资\n• 首期总投资 100亿元\n• 每个省办事处首期投资政府项目 5亿元"
    },
    {
        "title": "凝聚华人力量 · 共创美好未来",
        "body": "世界华人基金会\nWorld Chinese Foundation\n\n和平 · 友爱 · 发展 · 共赢"
    },
]

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for i, s in enumerate(slides):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    
    # 背景色矩形
    from pptx.util import Inches, Pt
    bg_color = "1a1a2e" if i % 2 == 0 else "16213e"
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.333), Inches(1.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = s["title"]
    p.font.size = Pt(48) if not s.get("sub") else Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Body
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.333), Inches(3.8))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = s["body"]
    p2.font.size = Pt(30) if not s.get("sub") else Pt(28)
    p2.alignment = PP_ALIGN.CENTER

output = "/tmp/world_chinese_foundation.pptx"
prs.save(output)
print(f"✅ PPTX saved: {output}")
print(f"Pages: {len(slides)}")
