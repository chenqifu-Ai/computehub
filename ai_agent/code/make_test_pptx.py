#!/usr/bin/env python3
"""生成测试 PPTX 用于视频生产管线测试"""
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.run(["pip3", "install", "python-pptx", "lxml"], check=True)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides_content = [
    ("小智影业 · 智能视频生产管线", "一键 PPT 转视频 · AI 语音合成 · 自动渲染"),
    ("核心技术栈", "• ComputeHub 分布式计算平台\n• edge-tts 语音合成引擎\n• FFmpeg 视频编码\n• Python 全自动管线"),
    ("视频生产流程", "1. 上传 PPT/文档 → 2. 逐页渲染为高清图片\n3. AI 语音合成（支持多种音色）\n4. 编码为视频段 → 5. 拼接成完整视频"),
    ("支持格式", "• PPTX / DOCX 一键转换\n• PDF 文档处理\n• 图片轮播\n• 纯文字脚本自动生成"),
    ("应用场景", "• 商业演示 → 宣传视频\n• 培训课件 → 教学视频\n• 项目方案 → 汇报视频\n• 产品介绍 → 营销视频"),
    ("部署架构", "Gateway（API调度）\n  └─ Worker（并行执行）\n       ├─ 渲染节点（PIL+FFmpeg）\n       └─ 语音节点（edge-tts）"),
    ("Gallery 作品广场", "• 自动归档已生成视频\n• 支持在线预览和下载\n• 版本管理和分类标签\n• 后续支持评论区互动"),
    ("视频参数", "• 分辨率：1920×1080 全高清\n• 帧率：30fps\n• 编码：H.264 + MP3\n• 时长：按内容自动适配"),
    ("模板风格", "• Business — 商务蓝白配色\n• Clean — 简洁黑白风格\n• Minimal — 极简展示\n• 可自定义配色和字体"),
    ("语音合成", "• 云希（男声）— 稳重专业\n• 晓晓（女声）— 甜美亲切\n• 云扬（男声）— 活力朝气\n• 晓辰（女声）— 温柔知性"),
    ("性能指标", "• 单页渲染：<2秒\n• 语音合成：<3秒/页\n• 编码速度：约 3x 实时\n• 10页PPT：<2分钟完成"),
    ("未来规划", "• 自定义背景音乐\n• 高级字幕特效\n• AI 智能脚本生成\n• 批量处理队列"),
    ("感谢观看", "小智影业 · AI 驱动的智能视频生产\nComputeHub 赋能内容创作"),
]

for title, body in slides_content:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    from pptx.util import Inches, Pt, Emu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = None  # inherit
    p.alignment = PP_ALIGN.CENTER
    
    # Body
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9.333), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body
    p2.font.size = Pt(28)
    p2.alignment = PP_ALIGN.CENTER

output_path = "/root/.openclaw/workspace/test_video_pipeline.pptx"
prs.save(output_path)
print(f"✅ 测试 PPTX 已生成: {output_path}")
print(f"   共 {len(slides_content)} 页")
