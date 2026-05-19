#!/usr/bin/env python3
"""
AI智能招商演示主机商业计划书 → PPTX生成 + 视频管线一键执行
"""
import sys
import os

# ─── Step 1: 在 ECS 上生成 PPTX ──────────────────────────────

PPTX_SCRIPT = r'''
#!/usr/bin/env python3
"""生成：AI智能招商演示主机商业计划书（精简版）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLUE_DARK  = RGBColor(15, 30, 60)
BLUE_MID   = RGBColor(25, 50, 100)
BLUE_ACCENT = RGBColor(59, 130, 246)
GOLD       = RGBColor(218, 165, 32)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
SOFT_BLUE  = RGBColor(147, 197, 253)
BG_BLUE    = RGBColor(20, 40, 80)

def add_bg(slide, top=BLUE_DARK, bottom=BLUE_MID):
    """渐变背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = top

def add_accent_bar(slide):
    """顶部装饰条"""
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1, 0, 0, W, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_ACCENT
    shape.line.fill.background()

def add_footer(slide, text="厦门京云科技 · AI智能招商平台"):
    """底部文字"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = SOFT_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_title_text(slide, title, subtitle=None):
    """大标题 + 副标题"""
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER

def add_bullet_slide(slide, title, items, sub_items=None):
    """要点页"""
    # 标题
    txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    # 装饰条
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_ACCENT
    bar.line.fill.background()
    # 列表
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(4.8))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i > 0:
            tf2.add_paragraph()
        p = tf2.paragraphs[-1]
        p.text = f"● {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(10)
        if sub_items and i < len(sub_items) and sub_items[i]:
            tf2.add_paragraph()
            sp = tf2.paragraphs[-1]
            sp.text = f"    {sub_items[i]}"
            sp.font.size = Pt(16)
            sp.font.color.rgb = SOFT_BLUE
            sp.space_after = Pt(6)

def save_slide(slide):
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide)

# ═══════════════════════════════════════════════════════════════
# PAGE 1: 封面
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
save_slide(slide)
add_title_text(slide,
    "🚀 AI智能招商演示主机",
    "会说话的智能招商盒子 —— 张嘴就能讲全案")
# 公司名
txCo = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(7.3), Inches(0.6))
tf = txCo.text_frame
p = tf.paragraphs[0]
p.text = "厦门京云科技有限公司  |  2026"
p.font.size = Pt(18)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════
# PAGE 2: 痛点 · 传统招商之痛
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "🎯 招商之痛 · 传统模式四大困境", [
    "新人不会讲 —— 招商人员流动大，培训周期长",
    "标准不统一 —— 不同人讲不同样，品牌形象打折",
    "设备成本高 —— 投影 + 笔记本 + 专人操作",
    "互动体验差 —— 单向PPT播讲，无法实时应答互动",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 3: 解决方案
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "💡 AI智能招商演示主机 · 全新解决方案", [
    "一键自动讲解 —— 唤醒即讲，从头到尾专业演示",
    "语音交互控制 —— 说翻页 跳转 放大，零门槛操作",
    "内容统一管理 —— 后台一键更新，所有设备同步",
    "智能问答应答 —— 现场提问即刻回答，政策费用全掌握",
], sub_items=[
    "AI智能体自动全流程讲解，新人也专业",
    "离线/在线双模语音识别，精准率>95%",
    "多端同步，总部管控，版本统一",
    "RAG架构+LLM，知识库实时检索",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 4: 产品功能架构
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "🏗️ 五大功能模块", [
    "AI智能体模块 —— 语音识别 · 自动讲解 · 智能问答",
    "大屏交互模块 —— 语音翻页 · 视频播放 · 图文缩放",
    "全案管理模块 —— 品牌 · 产品 · 模式 · 投资 · 扶持",
    "数据统计模块 —— 高频问题 · 停留分析 · 转化监测",
    "管理后台模块 —— 内容发布 · 设备管理 · 权限控制",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 5: 核心优势
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "⚡ 六大核心优势", [
    "零培训上岗 —— 无需背诵，张嘴即讲，人人都是金牌招商",
    "7×24小时待命 —— 无人值守自动演示，访客随时观看",
    "内容安全加密 —— AES-256 + DRM，核心商业数据无忧",
    "硬件成本极低 —— 1台电视盒子替代整套投影方案",
    "快速迭代更新 —— 后台修改即刻生效，无版本滞后",
    "多场景适配 —— 招商会 · 门店 · 展会 · 线上路演",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 6: 商业模式
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_title_text(slide, "💼 商业模式 · 三维营收体系")
# 三个卡片
cards = [
    ("硬件销售", "¥2,980/台", "招商演示主机\n含一年平台服务"),
    ("SaaS年费", "¥4,800/年", "内容管理平台\n数据统计分析"),
    ("内容定制", "¥2万起", "招商全案定制\n视频/图文制作"),
]
for idx, (name, price, desc) in enumerate(cards):
    x = Inches(1.2 + idx * 3.8)
    y = Inches(4.0)
    box = slide.shapes.add_textbox(x, y, Inches(3.3), Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE_ACCENT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # 价格
    tf.add_paragraph()
    p2 = tf.paragraphs[-1]
    p2.text = price
    p2.font.size = Pt(32)
    p2.font.color.rgb = GOLD
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    # 描述
    tf.add_paragraph()
    p3 = tf.paragraphs[-1]
    p3.text = desc
    p3.font.size = Pt(15)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════
# PAGE 7: 市场分析
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "📊 市场前景与目标", [
    "目标客户 —— 连锁品牌方 · 招商加盟企业 · 区域代理商",
    "市场规模 —— 全国连锁品牌超50万家，每年新增10万+",
    "客户痛点 —— 90%品牌仍用传统PPT+人工模式招商",
    "三年目标 —— 服务5000家品牌，覆盖30%连锁招商市场",
    "营收预测 —— 第一年500万 → 第二年2000万 → 第三年5000万",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 8: 技术架构
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "🔧 技术架构 · 分层微服务", [
    "终端层 —— Android盒子 / 工控主机 · HDMI电视/投影",
    "交互层 —— 离线语音引擎 · 在线NLU · 语音合成(TTS)",
    "业务层 —— AI智能体 · 大屏交互 · 全案管理 · 数据",
    "AI引擎层 —— 语音识别 · RAG检索 · 大语言模型",
    "数据层 —— MySQL+Redis+MinIO · AES-256加密",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 9: 竞争优势
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "🏆 竞争壁垒", [
    "技术壁垒 —— 离线语音+在线AI双引擎，软硬一体集成",
    "数据壁垒 —— 持续积累招商问答库，越用越聪明",
    "先发优势 —— 国内首个专注招商场景的AI智能硬件方案",
    "生态壁垒 —— 内容管理平台绑定品牌方，迁移成本高",
    "成本壁垒 —— 电视盒子方案，价格仅为传统方案1/5",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 10: 团队与案例
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_bullet_slide(slide, "👥 研发团队 · 落地案例", [
    "核心团队 —— 10年AI+IoT经验，曾服务华为/海康/美团",
    "研发实力 —— 全栈自研：语音/视觉/大模型/嵌入式",
    "测试数据 —— 语音识别准确率96.2%，响应延迟<400ms",
    "落地案例 —— 已与3家连锁品牌达成试点合作",
    "研发周期 —— 180天交付，含硬件/软件/AI全套方案",
])

# ═══════════════════════════════════════════════════════════════
# PAGE 11: 融资计划
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_title_text(slide, "💰 融资计划 · 携手共赢")
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    "融资金额：¥500万  |  出让股份：10%",
    "资金用途：产品迭代(60%) · 市场推广(25%) · 团队扩建(15%)",
    "预期回本：12-18个月  |  IRR > 180%",
]
for i, item in enumerate(items):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"◆ {item}"
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_GRAY
    p.space_after = Pt(12)

# ═══════════════════════════════════════════════════════════════
# PAGE 12: 结尾
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
save_slide(slide)
add_title_text(slide,
    "🤝 让招商变得更简单",
    "厦门京云科技有限公司  · 诚邀合作伙伴")
txBox = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9.3), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "www.jingyun.tech  |  contact@jingyun.tech  |  400-888-0000"
p.font.size = Pt(18)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# ── 保存 ───────────────────────────────────────────────────
output = "/home/computehub/ai_zhaoshang_bp.pptx"
prs.save(output)
print(f"✅ PPTX saved: {output} ({os.path.getsize(output)/1024:.0f}KB)")
'''

import subprocess
import sys

# scp 脚本到 ECS
script_path = "/tmp/gen_pptx.py"
with open(script_path, "w") as f:
    f.write(PPTX_SCRIPT)

# 上传并执行
scp_cmd = [
    "scp", "-i", "/root/.ssh/id_ed25519_computehub",
    "-o", "StrictHostKeyChecking=no",
    script_path,
    "computehub@36.250.122.43:/tmp/gen_pptx.py"
]
print("📤 Uploading PPTX generator to ECS...")
subprocess.run(scp_cmd, check=True)

ssh_cmd = [
    "ssh", "-i", "/root/.ssh/id_ed25519_computehub",
    "-o", "StrictHostKeyChecking=no",
    "computehub@36.250.122.43",
    "python3 /tmp/gen_pptx.py"
]
print("🚀 Generating PPTX on ECS...")
result = subprocess.run(ssh_cmd, capture_output=True, text=True, timeout=30)
print(result.stdout)
if result.stderr:
    print("stderr:", result.stderr[:500])

# ─── Step 2: 执行视频管线 ──────────────────────────────────
print()
print("=" * 60)
print("🎬 Now running video pipeline on ECS...")
print("=" * 60)

pipeline_cmd = [
    "ssh", "-i", "/root/.ssh/id_ed25519_computehub",
    "-o", "StrictHostKeyChecking=no",
    "computehub@36.250.122.43",
    "cd /home/computehub/scripts/video && "
    "python3 video_pipeline.py "
    "--doc /home/computehub/ai_zhaoshang_bp.pptx "
    "--output /home/computehub/gallery/ai_zhaoshang_bp.mp4 "
    "--template tech_blue "
    "--voice xiaochen "
    "--title 'AI智能招商演示主机商业计划书' "
    "--bgm '' "
    "2>&1"
]

print("⏳ Pipeline running (can take 3-5 min)...")
result2 = subprocess.run(pipeline_cmd, capture_output=True, text=True, timeout=600)
print(result2.stdout)
if result2.stderr:
    print("stderr:", result2.stderr[:1000])

print()
print("🎉 Done!")
