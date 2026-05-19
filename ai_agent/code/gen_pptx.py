#!/usr/bin/env python3
"""
Generate PPTX: AI智能招商演示主机商业计划书 (精简版)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BLUE_D   = RGBColor(15,30,60)
BLUE_M   = RGBColor(25,50,100)
BLUE_A   = RGBColor(59,130,246)
GOLD     = RGBColor(218,165,32)
WHITE    = RGBColor(255,255,255)
LGRAY    = RGBColor(226,232,240)
SBLUE    = RGBColor(147,197,253)

def save_slide(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE_D
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_A
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "厦门京云科技 · AI智能招商平台"
    p.font.size = Pt(11)
    p.font.color.rgb = SBLUE
    p.alignment = PP_ALIGN.CENTER

def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(44); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle; p2.font.size = Pt(22); p2.font.color.rgb = LGRAY; p2.alignment = PP_ALIGN.CENTER

def add_bullets(slide, title, items):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(36); p.font.color.rgb = WHITE; p.font.bold = True
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE_A; bar.line.fill.background()
    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(4.8))
    tf = tb2.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i > 0: tf.add_paragraph()
        p = tf.paragraphs[-1]; p.text = "\u25cf " + item; p.font.size = Pt(20); p.font.color.rgb = LGRAY; p.space_after = Pt(10)

# ════ P1: 封面 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_title(s, "AI智能招商演示主机", "会说话的智能招商盒子 - 张嘴就能讲全案")
tb = s.shapes.add_textbox(Inches(3), Inches(5.2), Inches(7.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.text = "厦门京云科技有限公司  |  2026"; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER

# ════ P2: 痛点 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "传统招商四大困境", [
    "新人不会讲 - 招商人员流动大,培训周期长",
    "标准不统一 - 不同人讲不同样,品牌形象打折",
    "设备成本高 - 投影+笔记本+专人操作",
    "互动体验差 - 单向PPT播讲,无法实时应答",
])

# ════ P3: 解决方案 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "AI智能招商演示主机 - 全新方案", [
    "一键自动讲解 - AI唤醒即讲,从头到尾专业演示",
    "语音交互控制 - 说翻页跳转放大,零门槛操作",
    "内容统一管理 - 后台一键更新,所有设备同步",
    "智能问答交互 - 现场提问即刻回答,政策费用全掌握",
])

# ════ P4: 五大模块 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "五大功能模块", [
    "AI智能体 - 语音识别/自动讲解/智能问答",
    "大屏交互 - 语音翻页/视频播放/图文缩放",
    "全案管理 - 品牌/产品/模式/投资/扶持",
    "数据统计 - 高频问题/停留分析/转化监测",
    "管理后台 - 内容发布/设备管理/权限控制",
])

# ════ P5: 核心优势 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "六大核心优势", [
    "零培训上岗 - 无需背诵,张嘴即讲,人人金牌招商",
    "7x24待命 - 无人值守自动演示,访客随时观看",
    "内容安全 - AES-256+DRM,核心数据无忧",
    "极致成本 - 1台电视盒子替代整套投影方案",
    "快速迭代 - 后台修改即刻生效,无版本滞后",
    "多场景适配 - 招商会/门店/展会/线上路演",
])

# ════ P6: 商业模式 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_title(s, "商业模式 - 三维营收体系")
cards = [("硬件销售", "2,980/台", "招商演示主机\n含一年平台服务"),
         ("SaaS年费", "4,800/年", "内容管理平台\n数据统计分析"),
         ("内容定制", "2万起", "招商全案定制\n视频/图文制作")]
for idx,(name,price,desc) in enumerate(cards):
    x = Inches(1.2+idx*3.8); y = Inches(4.0)
    tb = s.shapes.add_textbox(x,y,Inches(3.3),Inches(2.8)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(24); p.font.color.rgb = BLUE_A; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.add_paragraph(); p2 = tf.paragraphs[-1]; p2.text = price; p2.font.size = Pt(32); p2.font.color.rgb = GOLD; p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER
    tf.add_paragraph(); p3 = tf.paragraphs[-1]; p3.text = desc; p3.font.size = Pt(15); p3.font.color.rgb = LGRAY; p3.alignment = PP_ALIGN.CENTER

# ════ P7: 市场分析 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "市场前景与目标", [
    "目标客户 - 连锁品牌方/招商加盟企业/区域代理商",
    "市场规模 - 全国连锁品牌超50万家,年新增10万+",
    "痛点痛点 - 90%品牌仍用传统PPT+人工模式招商",
    "三年目标 - 服务5000家品牌,覆盖30%连锁招商市场",
    "营收预测 - 第一年500万 第二年2000万 第三年5000万",
])

# ════ P8: 技术架构 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "技术架构 - 分层微服务", [
    "终端层 - Android盒子/工控主机 HDMI电视/投影",
    "交互层 - 离线语音引擎 在线NLU 语音合成TTS",
    "业务层 - AI智能体 大屏交互 全案管理 数据统计",
    "AI引擎层 - 语音识别 RAG检索 大语言模型",
    "数据层 - MySQL+Redis+MinIO AES-256加密",
])

# ════ P9: 竞争壁垒 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "竞争壁垒", [
    "技术壁垒 - 离线语音+在线AI双引擎,软硬一体集成",
    "数据壁垒 - 持续积累招商问答库,越用越聪明",
    "先发优势 - 国内首个专注招商场景的AI智能硬件方案",
    "生态壁垒 - 内容平台绑定品牌方,迁移成本高",
    "成本壁垒 - 电视盒子方案,价格仅为传统方案1/5",
])

# ════ P10: 团队案例 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_bullets(s, "研发团队与落地案例", [
    "核心团队 - 10年AI+IoT经验,曾服务华为/海康/美团",
    "研发实力 - 全栈自研:语音/视觉/大模型/嵌入式",
    "测试数据 - 语音识别准确率96.2%,响应延迟400ms以内",
    "落地案例 - 已与3家连锁品牌达成试点合作",
    "研发周期 - 180天交付,含硬件/软件/AI全套方案",
])

# ════ P11: 融资计划 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_title(s, "融资计划 - 携手共赢")
tb = s.shapes.add_textbox(Inches(1.5),Inches(3.8),Inches(10.3),Inches(3.0)); tf = tb.text_frame; tf.word_wrap = True
items = ["融资金额:500万  出让股份:10%",
         "资金用途:产品迭代60% 市场推广25% 团队扩建15%",
         "预期回本:12-18个月  IRR超过180%"]
for i,item in enumerate(items):
    if i: tf.add_paragraph()
    p = tf.paragraphs[-1]; p.text = "\u25c6 "+item; p.font.size = Pt(22); p.font.color.rgb = LGRAY; p.space_after = Pt(12)

# ════ P12: 结尾 ════
s = prs.slides.add_slide(prs.slide_layouts[6]); save_slide(s)
add_title(s, "让招商变得更简单", "厦门京云科技有限公司 - 诚邀合作伙伴")
tb = s.shapes.add_textbox(Inches(2),Inches(5.0),Inches(9.3),Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "www.jingyun.tech | contact@jingyun.tech | 400-888-0000"; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER

out = "/home/computehub/ai_zhaoshang_bp.pptx"
prs.save(out)
print(f"OK pptx={out} size={os.path.getsize(out)}")
