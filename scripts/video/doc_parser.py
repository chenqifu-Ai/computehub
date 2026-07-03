#!/usr/bin/env python3
"""
📄 ComputeHub 文档解析器
支持: PPTX / DOCX / PDF

输出统一的数据结构:
  {
    "title": str,
    "pages": [{"page_num": int, "text": str, "images": [str], "has_image": bool}, ...],
    "total_pages": int
  }

坑修复记录 (2026-05-16):
  #1 moviepy.editor → from moviepy import ImageClip 已迁移
  #2 坐标顺序修正
  #7 ffprobe float 提取改用 re.search()
"""
import os
import re
import json
import tempfile
import subprocess
import shutil
from pathlib import Path

# ── 可选依赖 ──────────────────────────────────────────────────
# PPTX
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# DOCX
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PDF
try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# 图片转换
try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


# ── 统一输出 ──────────────────────────────────────────────────

ParsedDoc = dict  # 符合上述 schema


def parse(file_path: str) -> ParsedDoc:
    """智能解析：根据扩展名自动选择解析器"""
    ext = Path(file_path).suffix.lower()

    if ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"):
        return parse_image(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def parse_pptx(file_path: str) -> ParsedDoc:
    """解析 PPTX 文件"""
    if not HAS_PPTX:
        raise ImportError("需要 python-pptx: pip install python-pptx")

    prs = Presentation(file_path)
    pages = []

    for i, slide in enumerate(prs.slides):
        texts = []
        images = []

        for shape in slide.shapes:
            # 提取文本
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

            # 提取图片
            if shape.shape_type == 13:  # Picture
                try:
                    img_data = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1]
                    img_path = f"_page{i}_img{len(images)}.{ext}"
                    os.makedirs("/tmp/video_parser", exist_ok=True)
                    full_path = os.path.join("/tmp/video_parser", img_path)
                    with open(full_path, "wb") as f:
                        f.write(img_data)
                    images.append(full_path)
                except Exception:
                    pass

        # 标题 = 第一行或第一个大号文本
        title = texts[0] if texts else f"第{i+1}页"
        page_text = "\n".join(texts)

        pages.append({
            "page_num": i + 1,
            "text": page_text,
            "title": title,
            "images": images,
            "has_image": len(images) > 0,
        })

    return {
        "title": Path(file_path).stem,
        "pages": pages,
        "total_pages": len(pages),
    }


def parse_docx(file_path: str) -> ParsedDoc:
    """解析 DOCX 文件"""
    if not HAS_DOCX:
        raise ImportError("需要 python-docx: pip install python-docx")

    doc = Document(file_path)

    # 按段落分段，每 ~10 段作为一页
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                ext = rel.target_ref.split(".")[-1]
                img_path = f"_img_{len(images)}.{ext}"
                os.makedirs("/tmp/video_parser", exist_ok=True)
                full_path = os.path.join("/tmp/video_parser", img_path)
                with open(full_path, "wb") as f:
                    f.write(img_data)
                images.append(full_path)
            except Exception:
                pass

    pages = []
    page_size = 10
    for i in range(0, len(paragraphs), page_size):
        chunk = paragraphs[i:i + page_size]
        title = chunk[0] if chunk else f"第{len(pages)+1}页"
        pages.append({
            "page_num": len(pages) + 1,
            "text": "\n".join(chunk),
            "title": title,
            "images": images if len(pages) == 0 else [],
            "has_image": len(images) > 0 and len(pages) == 0,
        })

    return {
        "title": Path(file_path).stem,
        "pages": pages,
        "total_pages": len(pages),
    }


def parse_pdf(file_path: str) -> ParsedDoc:
    """解析 PDF 文件"""
    if not HAS_PDF:
        raise ImportError("需要 pdfplumber: pip install pdfplumber")

    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            title = lines[0] if lines else f"第{i+1}页"

            # 转图片（如果 PIL 可用）
            img_paths = []
            if HAS_PIL:
                try:
                    # 用 pdf2image 或直接转
                    img = page.to_image(resolution=150)
                    img_path = f"_page{i}_pdf.png"
                    os.makedirs("/tmp/video_parser", exist_ok=True)
                    full_path = os.path.join("/tmp/video_parser", img_path)
                    img.save(full_path)
                    img_paths.append(full_path)
                except Exception:
                    pass

            pages.append({
                "page_num": i + 1,
                "text": text,
                "title": title,
                "images": img_paths,
                "has_image": len(img_paths) > 0,
            })

    return {
        "title": Path(file_path).stem,
        "pages": pages,
        "total_pages": len(pages),
    }


def parse_image(file_path: str) -> ParsedDoc:
    """单张图片 → 单页文档"""
    if not HAS_PIL:
        raise ImportError("需要 Pillow: pip install Pillow")

    return {
        "title": Path(file_path).stem,
        "pages": [{
            "page_num": 1,
            "text": Path(file_path).stem,
            "title": Path(file_path).stem,
            "images": [file_path],
            "has_image": True,
        }],
        "total_pages": 1,
    }


# ── 工具 ──────────────────────────────────────────────────────

def render_page_as_image(page: dict, output_path: str,
                          width: int = 1920, height: int = 1080,
                          template: str = "business"):
    """将解析后的页面渲染为 PIL 图片

    template: "business" | "clean" | "minimal"
    """
    if not HAS_PIL:
        raise ImportError("需要 Pillow")

    img = Image.new("RGB", (width, height), color=get_bg_color(template))
    draw = ImageDraw.Draw(img)

    # 找到中文字体
    font_path = find_chinese_font()
    font_large = ImageFont.truetype(font_path, 48) if font_path else ImageFont.load_default()
    font_normal = ImageFont.truetype(font_path, 36) if font_path else ImageFont.load_default()
    font_small = ImageFont.truetype(font_path, 24) if font_path else ImageFont.load_default()

    # 画标题
    title = page.get("title", "")
    if title:
        # 标题栏背景
        draw.rectangle([(0, 0), (width, 120)], fill=get_accent_color(template))
        # 标题文字 — 使用 textbbox 确定位置
        bbox = draw.textbbox((0, 0), title, font=font_large)
        tx = (width - (bbox[2] - bbox[0])) // 2
        ty = (120 - (bbox[3] - bbox[1])) // 2
        draw.text((tx, ty), title, fill="white", font=font_large)

    # 画页码
    page_num = page.get("page_num", 1)
    total = page.get("total_pages", 1)
    if total > 1:
        footer_text = f"{page_num} / {total}"
        bbox = draw.textbbox((0, 0), footer_text, font=font_small)
        draw.text(
            (width - (bbox[2] - bbox[0]) - 30, height - 40),
            footer_text,
            fill="gray",
            font=font_small,
        )

    # 画正文
    text = page.get("text", "")
    # 去掉标题行（已单独渲染）
    lines = text.split("\n")
    if lines and lines[0] == title:
        lines = lines[1:]

    y = 160
    margin = 80
    max_width = width - 2 * margin
    line_height = 48

    for line in lines:
        if not line.strip():
            y += line_height // 2
            continue
        # 自动换行
        words = list(line)
        current_line = ""
        for ch in words:
            test_line = current_line + ch
            bbox = draw.textbbox((0, 0), test_line, font=font_normal)
            if (bbox[2] - bbox[0]) > max_width:
                draw.text((margin, y), current_line, fill=get_text_color(template), font=font_normal)
                y += line_height
                current_line = ch
            else:
                current_line = test_line
        if current_line:
            draw.text((margin, y), current_line, fill=get_text_color(template), font=font_normal)
            y += line_height

        if y > height - 80:
            break

    # 如果有图片，在右下角叠加
    if page.get("has_image") and page.get("images"):
        img_path = page["images"][0]
        if os.path.exists(img_path):
            try:
                overlay = Image.open(img_path)
                # 缩放到 300x300 以内
                overlay.thumbnail((300, 300))
                # 右下角
                ix = width - overlay.width - 40
                iy = height - overlay.height - 80
                img.paste(overlay, (ix, iy), overlay if overlay.mode == "RGBA" else None)
            except Exception:
                pass

    img.save(output_path, "JPEG", quality=85)
    return output_path


# ── 模板颜色 ──────────────────────────────────────────────────

TEMPLATES = {
    "business": {
        "bg": (15, 23, 42),        # 深蓝黑
        "accent": (37, 99, 235),   # 蓝色
        "text": (226, 232, 240),   # 浅灰白
    },
    "clean": {
        "bg": (255, 255, 255),     # 纯白
        "accent": (15, 23, 42),    # 深色
        "text": (30, 41, 59),      # 深灰
    },
    "minimal": {
        "bg": (248, 250, 252),     # 浅灰
        "accent": (100, 116, 139), # 灰蓝
        "text": (15, 23, 42),      # 深色
    },
}


def get_bg_color(template: str):
    return TEMPLATES.get(template, TEMPLATES["business"])["bg"]


def get_accent_color(template: str):
    return TEMPLATES.get(template, TEMPLATES["business"])["accent"]


def get_text_color(template: str):
    return TEMPLATES.get(template, TEMPLATES["business"])["text"]


def find_chinese_font() -> str:
    """找到中文字体"""
    font_candidates = [
        "/system/fonts/MiSansVF.ttf",
        "/system/fonts/DroidSans.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for fp in font_candidates:
        if os.path.exists(fp):
            return fp
    # 搜索
    for root_dir in ["/system/fonts", "/usr/share/fonts"]:
        if os.path.isdir(root_dir):
            for dirpath, _, filenames in os.walk(root_dir):
                for fn in filenames:
                    if "misans" in fn.lower() and fn.endswith((".ttf", ".otf")):
                        return os.path.join(dirpath, fn)
    return "/system/fonts/DroidSans.ttf"


# ── 测试 ──────────────────────────────────────────────────────

def test_render():
    """测试渲染"""
    test_page = {
        "page_num": 1,
        "title": "测试标题 — ComputeHub 视频生成",
        "text": "这是第一行正文。\n这是第二行，包含一些更长的文字来测试自动换行功能。\n\n第三段内容。",
        "images": [],
        "has_image": False,
    }
    out = "/tmp/test_render.jpg"
    render_page_as_image(test_page, out)
    print(f"✅ 渲染测试: {out}")
    return out


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        result = parse(sys.argv[1])
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        # 自测
        test_render()
        print("✅ 文档解析器自测通过")
